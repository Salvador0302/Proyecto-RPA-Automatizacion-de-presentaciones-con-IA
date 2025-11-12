#!/usr/bin/env python3
# scripts/latex_to_pptx.py
"""
Convierte archivos LaTeX/PDF en presentaciones PowerPoint
"""
import sys
import os
import subprocess
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches


class LatexToPptxConverter:
    """Clase para convertir archivos LaTeX/PDF a PowerPoint"""
    
    def __init__(self, dpi=150):
        """
        Inicializa el convertidor
        
        Args:
            dpi: Resolución para convertir PDF a imágenes
        """
        self.dpi = dpi
    
    def convert(self, input_file, output_file):
        """
        Convierte un archivo LaTeX o PDF a PowerPoint
        
        Args:
            input_file: Ruta del archivo LaTeX (.tex) o PDF (.pdf)
            output_file: Ruta del archivo PPTX de salida
        """
        if not os.path.exists(input_file):
            raise FileNotFoundError(f"Archivo no encontrado: {input_file}")
        
        # Si es un archivo LaTeX, primero compilarlo a PDF
        if input_file.endswith('.tex'):
            pdf_file = self._compile_latex(input_file)
        elif input_file.endswith('.pdf'):
            pdf_file = input_file
        else:
            raise ValueError(f"Formato no soportado: {input_file}. Use .tex o .pdf")
        
        # Convertir PDF a PPTX
        self._pdf_to_pptx(pdf_file, output_file)
        
        # Limpiar archivo PDF temporal si se compiló desde LaTeX
        if input_file.endswith('.tex') and pdf_file != input_file:
            try:
                os.remove(pdf_file)
            except:
                pass
    
    def _compile_latex(self, tex_file):
        """
        Compila un archivo LaTeX a PDF
        
        Args:
            tex_file: Ruta del archivo .tex
            
        Returns:
            str: Ruta del archivo PDF generado
        """
        # Obtener el directorio y nombre base del archivo
        tex_dir = os.path.dirname(tex_file) or '.'
        tex_basename = os.path.splitext(os.path.basename(tex_file))[0]
        pdf_file = os.path.join(tex_dir, f"{tex_basename}.pdf")
        
        try:
            # Intentar compilar con pdflatex
            result = subprocess.run(
                ['pdflatex', '-interaction=nonstopmode', f'-output-directory={tex_dir}', tex_file],
                capture_output=True,
                text=True,
                timeout=30
            )
            
            if not os.path.exists(pdf_file):
                raise Exception(f"No se pudo generar el PDF. pdflatex output: {result.stderr}")
            
            return pdf_file
            
        except FileNotFoundError:
            raise Exception(
                "pdflatex no está instalado o no está en el PATH. "
                "Por favor, instala una distribución de LaTeX (MiKTeX, TeX Live, etc.)"
            )
        except subprocess.TimeoutExpired:
            raise Exception("La compilación de LaTeX tardó demasiado tiempo")
        except Exception as e:
            raise Exception(f"Error al compilar LaTeX: {str(e)}")
    
    def _pdf_to_pptx(self, pdf_path, pptx_path):
        """
        Convierte un PDF a PowerPoint
        
        Args:
            pdf_path: Ruta del archivo PDF
            pptx_path: Ruta del archivo PPTX de salida
        """
        if not os.path.exists(pdf_path):
            raise FileNotFoundError(f"PDF no encontrado: {pdf_path}")
        
        # Convertir páginas del PDF a imágenes
        pages = convert_from_path(pdf_path, dpi=self.dpi)
        
        # Crear presentación
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]  # layout vacío
        
        for i, img in enumerate(pages):
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Guardar imagen temporalmente
            img_path = f"_tmp_slide_{i}.png"
            img.save(img_path, "PNG")
            
            # Insertar imagen ocupando todo el slide
            slide.shapes.add_picture(
                img_path,
                Inches(0),
                Inches(0),
                width=prs.slide_width,
                height=prs.slide_height
            )
            
            # Limpiar imagen temporal
            try:
                os.remove(img_path)
            except:
                pass
        
        # Crear directorio si no existe
        os.makedirs(os.path.dirname(pptx_path) or '.', exist_ok=True)
        
        prs.save(pptx_path)
        print(f"Presentación guardada: {pptx_path}")


# Función legacy para compatibilidad
def pdf_to_pptx(pdf_path, pptx_path, dpi=150):
    """Función legacy - usar LatexToPptxConverter en su lugar"""
    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"PDF not found: {pdf_path}")
    pages = convert_from_path(pdf_path, dpi=dpi)  # lista de PIL images
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]  # layout vacío

    for i, img in enumerate(pages):
        slide = prs.slides.add_slide(blank_slide_layout)
        img_path = f"_tmp_slide_{i}.png"
        img.save(img_path, "PNG")
        # insertar imagen ocupando todo el slide
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
        os.remove(img_path)

    prs.save(pptx_path)
    print("Saved:", pptx_path)


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python scripts/latex_to_pptx.py input.pdf output.pptx")
        sys.exit(1)
    pdf_to_pptx(sys.argv[1], sys.argv[2])
