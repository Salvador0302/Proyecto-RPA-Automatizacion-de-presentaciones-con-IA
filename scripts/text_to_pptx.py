#!/usr/bin/env python3
# scripts/text_to_pptx.py
"""
Convierte archivos de texto estructurados en presentaciones PowerPoint
"""
import os
import re
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


class TextToPptxConverter:
    """Clase para convertir texto estructurado a PowerPoint"""
    
    def __init__(self, theme="modern_blue"):
        """
        Inicializa el convertidor
        
        Args:
            theme: Tema de colores ("modern_blue", "dark", "professional", "vibrant")
        """
        self.prs = None
        self.theme = theme
        self.colors = self._get_theme_colors(theme)
    
    def _get_theme_colors(self, theme):
        """Obtiene los colores según el tema seleccionado"""
        themes = {
            "modern_blue": {
                "primary": RGBColor(41, 128, 185),      # Azul moderno
                "secondary": RGBColor(52, 152, 219),    # Azul claro
                "accent": RGBColor(46, 204, 113),       # Verde
                "text": RGBColor(44, 62, 80),           # Gris oscuro
                "bg_title": RGBColor(236, 240, 241),    # Gris muy claro
                "bg_content": RGBColor(255, 255, 255)   # Blanco
            },
            "dark": {
                "primary": RGBColor(41, 128, 185),
                "secondary": RGBColor(142, 68, 173),    # Púrpura
                "accent": RGBColor(230, 126, 34),       # Naranja
                "text": RGBColor(236, 240, 241),        # Texto claro
                "bg_title": RGBColor(44, 62, 80),       # Fondo oscuro
                "bg_content": RGBColor(52, 73, 94)      # Fondo medio
            },
            "professional": {
                "primary": RGBColor(51, 51, 51),        # Gris oscuro
                "secondary": RGBColor(0, 114, 188),     # Azul corporativo
                "accent": RGBColor(255, 185, 0),        # Dorado
                "text": RGBColor(51, 51, 51),
                "bg_title": RGBColor(242, 242, 242),
                "bg_content": RGBColor(255, 255, 255)
            },
            "vibrant": {
                "primary": RGBColor(231, 76, 60),       # Rojo
                "secondary": RGBColor(155, 89, 182),    # Púrpura
                "accent": RGBColor(241, 196, 15),       # Amarillo
                "text": RGBColor(44, 62, 80),
                "bg_title": RGBColor(236, 240, 241),
                "bg_content": RGBColor(255, 255, 255)
            }
        }
        return themes.get(theme, themes["modern_blue"])
    
    def convert(self, input_file, output_file):
        """
        Convierte un archivo de texto a PowerPoint
        
        Args:
            input_file: Ruta del archivo de texto de entrada
            output_file: Ruta del archivo PPTX de salida
        """
        if not os.path.exists(input_file):
            raise FileNotFoundError(f"Archivo no encontrado: {input_file}")
        
        # Leer el contenido del archivo
        with open(input_file, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # Parsear el contenido
        structure = self._parse_content(content)
        
        # Crear la presentación
        self._create_presentation(structure, output_file)
    
    def _parse_content(self, content):
        """
        Parsea el contenido de texto y lo convierte en una estructura
        
        Args:
            content: Contenido del archivo de texto
            
        Returns:
            list: Lista de diccionarios con estructura de slides
        """
        structure = []
        slides = re.split(r'SLIDE \d+:', content)
        
        for slide_content in slides:
            if not slide_content.strip():
                continue
            
            lines = [line.strip() for line in slide_content.split('\n') if line.strip()]
            if not lines:
                continue
            
            # La primera línea es el título
            title = lines[0]
            
            # Las siguientes líneas que empiezan con - son bullets
            bullets = []
            for line in lines[1:]:
                if line.startswith('-'):
                    bullets.append(line[1:].strip())
                elif line and not line.startswith('SLIDE'):
                    # Si no empieza con -, puede ser continuación del título o bullet
                    if not bullets:
                        title += " " + line
                    else:
                        bullets.append(line)
            
            structure.append({
                "title": title,
                "bullets": bullets
            })
        
        return structure
    
    def _create_presentation(self, structure, output_file):
        """
        Crea la presentación PowerPoint a partir de la estructura
        
        Args:
            structure: Lista de diccionarios con estructura de slides
            output_file: Ruta del archivo PPTX de salida
        """
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
        for i, slide_data in enumerate(structure):
            if i == 0:
                # Primera diapositiva como portada
                self._make_title_slide(
                    self.prs,
                    slide_data.get("title", ""),
                    slide_data.get("bullets", [])
                )
            else:
                # Diapositivas de contenido
                self._make_content_slide(
                    self.prs,
                    slide_data.get("title", ""),
                    slide_data.get("bullets", [])
                )
        
        # Crear directorio si no existe
        os.makedirs(os.path.dirname(output_file) or '.', exist_ok=True)
        
        self.prs.save(output_file)
        print(f"Presentación guardada: {output_file}")
    
    def _make_title_slide(self, prs, title, bullets):
        """Crea una diapositiva de título (portada)"""
        # Usar layout en blanco
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Fondo de la diapositiva
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors["primary"]
        
        # Título principal
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(1.5)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.word_wrap = True
        
        # Formato del título
        p = title_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Subtítulo/descripción si hay bullets
        if bullets:
            subtitle_top = Inches(4.5)
            subtitle_box = slide.shapes.add_textbox(left, subtitle_top, width, Inches(1))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = " | ".join(bullets[:3])  # Máximo 3 bullets
            
            p_sub = subtitle_frame.paragraphs[0]
            p_sub.alignment = PP_ALIGN.CENTER
            p_sub.font.size = Pt(20)
            p_sub.font.color.rgb = RGBColor(255, 255, 255)
        
        # Línea decorativa
        line_top = Inches(4.2)
        line = slide.shapes.add_shape(
            1,  # Rectángulo
            Inches(3.5), line_top,
            Inches(3), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.colors["accent"]
        line.line.fill.background()
    
    def _make_content_slide(self, prs, title, bullets):
        """Crea una diapositiva de contenido con diseño moderno"""
        # Usar layout en blanco
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Fondo blanco/claro
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors["bg_content"]
        
        # Barra superior decorativa
        header_shape = slide.shapes.add_shape(
            1,  # Rectángulo
            Inches(0), Inches(0),
            Inches(10), Inches(0.8)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = self.colors["primary"]
        header_shape.line.fill.background()
        
        # Título en la barra superior
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.word_wrap = True
        
        p_title = title_frame.paragraphs[0]
        p_title.font.size = Pt(32)
        p_title.font.bold = True
        p_title.font.color.rgb = RGBColor(255, 255, 255)
        
        # Área de contenido
        if bullets:
            content_left = Inches(1.2)
            content_top = Inches(1.5)
            content_width = Inches(8)
            content_height = Inches(5)
            
            content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            text_frame.vertical_anchor = 1  # Centro vertical
            
            for i, bullet in enumerate(bullets):
                if i > 0:
                    text_frame.add_paragraph()
                
                p = text_frame.paragraphs[i]
                p.text = bullet
                p.level = 0
                p.font.size = Pt(20)
                p.font.color.rgb = self.colors["text"]
                p.space_before = Pt(12)
                p.space_after = Pt(12)
                
                # Agregar viñeta personalizada
                p.font.bold = False
                
                # Icono de viñeta (bullet point)
                run = p.runs[0]
                run.text = f"● {bullet}"
                run.font.color.rgb = self.colors["secondary"]
        
        # Barra inferior decorativa
        footer_shape = slide.shapes.add_shape(
            1,  # Rectángulo
            Inches(0), Inches(7.2),
            Inches(10), Inches(0.3)
        )
        footer_shape.fill.solid()
        footer_shape.fill.fore_color.rgb = self.colors["accent"]
        footer_shape.line.fill.background()
    
    def _make_slide(self, prs, title, bullets):
        """Función legacy para compatibilidad"""
        self._make_content_slide(prs, title, bullets)


# Funciones legacy para compatibilidad
def make_slide(prs, title, bullets):
    """Función legacy - usar TextToPptxConverter en su lugar"""
    slide_layout = prs.slide_layouts[1]  # título y contenido
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    # limpia el primer párrafo generado por defecto
    body.clear()
    for b in bullets:
        p = body.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(18)


def create_pptx_from_structure(structure, out="output.pptx"):
    """Función legacy - usar TextToPptxConverter en su lugar"""
    prs = Presentation()
    for s in structure:
        make_slide(prs, s.get("title",""), s.get("bullets", []))
    prs.save(out)
    print("Saved:", out)


# ejemplo de uso
if __name__ == "__main__":
    sample = [
        {"title":"Slide 1", "bullets":["Punto A", "Punto B"]},
        {"title":"Slide 2", "bullets":["Otro punto"]}
    ]
    create_pptx_from_structure(sample, "example_from_text.pptx")
